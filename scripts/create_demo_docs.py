#!/usr/bin/env python3
"""Generate demo dialog DOCX and architecture PPTX for ShieldLend hackathon submission."""

import os
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt, Emu
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT_DIR = os.path.join(os.path.dirname(__file__), "..", "docs")
os.makedirs(OUT_DIR, exist_ok=True)


# ============================================================
# 1. DEMO DIALOG DOCX
# ============================================================
def create_demo_docx():
    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)

    # Title
    title = doc.add_heading("ShieldLend — 3-Minute Video Demo Script", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_paragraph(
        "Privacy-Native BTC Lending Protocol on Starknet\n"
        "Re{define} Hackathon | March 2026\n"
        "Tracks: Privacy + Bitcoin",
        style="Subtitle",
    ).alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_paragraph()

    # ---- Overview table ----
    doc.add_heading("Video Structure", level=1)
    table = doc.add_table(rows=7, cols=3, style="Light Shading Accent 1")
    headers = ["Segment", "Duration", "What to Show"]
    for i, h in enumerate(headers):
        table.rows[0].cells[i].text = h
    rows = [
        ("1. Intro / Hook", "0:00 – 0:20", "Title slide + problem statement"),
        ("2. Architecture", "0:20 – 0:50", "Architecture slide + contract structure"),
        ("3. Live Demo — Deposit", "0:50 – 1:20", "Connect wallet → deposit strkBTC"),
        ("4. Live Demo — Borrow", "1:20 – 1:50", "Borrow USDC against collateral"),
        ("5. Privacy Feature", "1:50 – 2:20", "Show shielded mode toggle + explain ZK flow"),
        ("6. Yield + Markets", "2:20 – 2:50", "Show yield tokenization + market stats"),
    ]
    for idx, (seg, dur, show) in enumerate(rows, 1):
        table.rows[idx].cells[0].text = seg
        table.rows[idx].cells[1].text = dur
        table.rows[idx].cells[2].text = show

    doc.add_paragraph()

    # ---- Detailed Script ----
    doc.add_heading("Detailed Script (Dialog + Actions)", level=1)

    # SEGMENT 1
    doc.add_heading("Segment 1 — Intro / Hook  (0:00 – 0:20)", level=2)
    p = doc.add_paragraph()
    p.add_run("[SLIDE: Title Slide — ShieldLend logo + tagline]\n\n").bold = True
    p.add_run("DIALOG:\n").bold = True
    doc.add_paragraph(
        '"Hi, I\'m 0xhaz and this is ShieldLend — a privacy-native BTC lending protocol '
        'built on Starknet.\n\n'
        'The problem: In DeFi lending today, every deposit, borrow, and liquidation is '
        'fully visible on-chain. Anyone can track your positions, front-run your trades, '
        'or target you for liquidation.\n\n'
        'ShieldLend solves this by making privacy a first-class feature. You can deposit BTC, '
        'borrow stablecoins, and earn yield — all with optional ZK-powered privacy."',
        style="Quote",
    )

    # SEGMENT 2
    doc.add_heading("Segment 2 — Architecture Overview  (0:20 – 0:50)", level=2)
    p = doc.add_paragraph()
    p.add_run("[SLIDE: Architecture Diagram]\n\n").bold = True
    p.add_run("DIALOG:\n").bold = True
    doc.add_paragraph(
        '"Let me walk you through the architecture.\n\n'
        'We built 16 Cairo smart contracts with over 3,200 lines of code, all passing 123 tests.\n\n'
        'At the core, we have isolated LendingPool markets — strkBTC/USDC, tBTC/USDC, and a '
        'strkBTC/tBTC eMode pool with 95% LTV for correlated BTC pairs.\n\n'
        'The privacy engine uses Pedersen commitments stored in a Merkle tree, a nullifier '
        'registry for double-spend prevention, and an on-chain ZK verifier — all native to Starknet.\n\n'
        'On top, we have yield tokenization that splits deposit receipts into Principal Tokens and '
        'Yield Tokens for advanced DeFi strategies, plus Dutch auction liquidations and flash loans."',
        style="Quote",
    )

    # SEGMENT 3
    doc.add_heading("Segment 3 — Live Demo: Deposit  (0:50 – 1:20)", level=2)
    p = doc.add_paragraph()
    p.add_run("[SCREEN: Browser showing ShieldLend frontend]\n\n").bold = True
    p.add_run("ACTIONS:\n").bold = True
    actions = [
        'Open the ShieldLend app in the browser',
        'Show the homepage — point out the live market stats (TVL, Active Markets, ZK-Ready badge)',
        'Click "Explore Markets" → show the 3 markets with live APYs and utilization rates',
        'Click on strkBTC/USDC market → show the market detail page with interest rate curve',
        'Navigate to "Lend" page',
        'Connect wallet (Argent X / Braavos)',
        'Select strkBTC/USDC market',
        'Enter deposit amount (e.g., 0.5 strkBTC)',
        'Click "Deposit" → approve + confirm transaction',
        'Show the updated balance — slTokens received',
    ]
    for a in actions:
        doc.add_paragraph(a, style="List Number")

    p = doc.add_paragraph()
    p.add_run("\nDIALOG:\n").bold = True
    doc.add_paragraph(
        '"Here\'s our live app deployed on Starknet Sepolia. You can see 3 active markets '
        'with real on-chain data — supply APY, borrow APY, and utilization rates all calculated '
        'from our two-slope interest rate model.\n\n'
        'Let me deposit some strkBTC as collateral. I connect my wallet, select the amount, '
        'and deposit. The pool takes my strkBTC and mints slTokens — these are my deposit '
        'receipt tokens that appreciate as borrowers pay interest."',
        style="Quote",
    )

    # SEGMENT 4
    doc.add_heading("Segment 4 — Live Demo: Borrow  (1:20 – 1:50)", level=2)
    p = doc.add_paragraph()
    p.add_run("[SCREEN: Borrow page]\n\n").bold = True
    p.add_run("ACTIONS:\n").bold = True
    actions = [
        'Navigate to "Borrow" page',
        'Select strkBTC/USDC market',
        'Show the collateral info — deposited strkBTC amount, health factor',
        'Enter borrow amount (e.g., 20,000 USDC)',
        'Point out the LTV ratio (75%) and health factor indicator',
        'Click "Borrow" → confirm transaction',
        'Show USDC received in wallet',
    ]
    for a in actions:
        doc.add_paragraph(a, style="List Number")

    p = doc.add_paragraph()
    p.add_run("\nDIALOG:\n").bold = True
    doc.add_paragraph(
        '"Now I can borrow against my collateral. The strkBTC/USDC pool has a 75% LTV ratio. '
        'The oracle gives us BTC prices, and the pool calculates how much I can borrow.\n\n'
        'I\'ll borrow 20,000 USDC. The pool verifies my collateral covers the loan, mints a '
        'DebtToken to track my debt, and sends USDC to my wallet. The interest rate adjusts '
        'dynamically based on pool utilization."',
        style="Quote",
    )

    # SEGMENT 5
    doc.add_heading("Segment 5 — Privacy Feature  (1:50 – 2:20)", level=2)
    p = doc.add_paragraph()
    p.add_run("[SCREEN: Lend page with Privacy toggle / Architecture slide]\n\n").bold = True
    p.add_run("ACTIONS:\n").bold = True
    actions = [
        'Show the Privacy Mode toggle on the Lend page',
        'Toggle it ON — explain what changes',
        'Point to the privacy indicator badge',
        '[SLIDE: Switch to Shielded Flow diagram]',
        'Walk through: commitment → proof → verify → nullify',
    ]
    for a in actions:
        doc.add_paragraph(a, style="List Number")

    p = doc.add_paragraph()
    p.add_run("\nDIALOG:\n").bold = True
    doc.add_paragraph(
        '"This is what makes ShieldLend different — privacy mode.\n\n'
        'When you toggle privacy on, instead of a normal deposit, the client generates a '
        'Pedersen commitment — hash of your amount plus a random blinding factor. This '
        'commitment goes into a Merkle tree on-chain via our CommitmentStore.\n\n'
        'When you borrow or withdraw, you reveal a nullifier — derived from your secret — '
        'which the NullifierRegistry checks to prevent double-spending. The ZK verifier '
        'confirms your proof without ever learning the actual amounts.\n\n'
        'The transparent and shielded modes use the SAME LendingPool contract but different '
        'entrypoints — so liquidity is shared but accounting is completely separate. This is '
        'all native to Starknet using STARKs, no external proving systems needed."',
        style="Quote",
    )

    # SEGMENT 6
    doc.add_heading("Segment 6 — Yield + Markets  (2:20 – 2:50)", level=2)
    p = doc.add_paragraph()
    p.add_run("[SCREEN: Yield page + Markets page]\n\n").bold = True
    p.add_run("ACTIONS:\n").bold = True
    actions = [
        'Navigate to "Yield" page',
        'Show the 3 yield tokenizer markets with maturity dates',
        'Explain PT (Principal Token) and YT (Yield Token) concept',
        'Navigate to "Portfolio" page — show positions overview',
        'Navigate back to Markets — show all 3 markets with live data',
    ]
    for a in actions:
        doc.add_paragraph(a, style="List Number")

    p = doc.add_paragraph()
    p.add_run("\nDIALOG:\n").bold = True
    doc.add_paragraph(
        '"We also built yield tokenization. You can split your slTokens — the deposit receipt — '
        'into Principal Tokens and Yield Tokens. PT gives you fixed-rate exposure, YT gives you '
        'leveraged yield exposure. After maturity, PT redeems 1:1 for slTokens and YT claims '
        'accrued interest.\n\n'
        'Plus we have flash loans with a 0.09% fee, Dutch auction liquidations, and eMode '
        'for correlated BTC pairs with up to 95% LTV."',
        style="Quote",
    )

    # CLOSING
    doc.add_heading("Closing  (2:50 – 3:00)", level=2)
    p = doc.add_paragraph()
    p.add_run("[SLIDE: Summary / Thank You slide]\n\n").bold = True
    p.add_run("DIALOG:\n").bold = True
    doc.add_paragraph(
        '"To recap — ShieldLend brings privacy-native BTC lending to Starknet. 16 contracts, '
        '123 tests, live on Sepolia, with ZK-powered privacy, yield tokenization, and a full '
        'production-ready frontend.\n\n'
        'Thank you for watching. The code is open source on GitHub."',
        style="Quote",
    )

    # ---- Tips ----
    doc.add_paragraph()
    doc.add_heading("Recording Tips", level=1)
    tips = [
        "Pre-record the wallet transactions to avoid waiting for block confirmations on camera",
        "Use screen recording software (OBS) at 1080p, record audio separately for clean sound",
        "Have the architecture slides ready to switch to during the privacy explanation",
        "Keep the browser zoomed in so text is readable in the video",
        "Practice the script 2-3 times to hit the 3-minute mark naturally",
    ]
    for t in tips:
        doc.add_paragraph(t, style="List Bullet")

    path = os.path.join(OUT_DIR, "ShieldLend_Demo_Script.docx")
    doc.save(path)
    print(f"Created: {path}")
    return path


# ============================================================
# 2. ARCHITECTURE PPTX
# ============================================================
def add_slide(prs, layout_idx, title_text, content_func):
    slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(slide_layout)
    # Dark background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = PptxRGB(0x0F, 0x17, 0x2A)
    content_func(slide, title_text)
    return slide


def set_title(slide, text, font_size=36):
    if slide.shapes.title:
        tf = slide.shapes.title.text_frame
        tf.text = text
        for para in tf.paragraphs:
            para.font.color.rgb = PptxRGB(0xFF, 0xFF, 0xFF)
            para.font.size = PptxPt(font_size)
            para.font.bold = True


def add_text_box(slide, left, top, width, height, text, font_size=14, color=PptxRGB(0xD1, 0xD5, 0xDB), bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(PptxInches(left), PptxInches(top), PptxInches(width), PptxInches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = PptxPt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_box(slide, left, top, width, height, text, fill_rgb, font_size=11, text_color=PptxRGB(0xFF, 0xFF, 0xFF)):
    """Add a rounded rectangle with text."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, PptxInches(left), PptxInches(top),
        PptxInches(width), PptxInches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = PptxPt(font_size)
    tf.paragraphs[0].font.color.rgb = text_color
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].space_before = PptxPt(4)
    shape.text_frame.paragraphs[0].space_after = PptxPt(4)
    return shape


def add_arrow(slide, start_left, start_top, end_left, end_top):
    """Add a connector arrow."""
    from pptx.enum.shapes import MSO_SHAPE
    connector = slide.shapes.add_connector(
        1,  # straight connector
        PptxInches(start_left), PptxInches(start_top),
        PptxInches(end_left), PptxInches(end_top)
    )
    connector.line.color.rgb = PptxRGB(0x6B, 0x72, 0x80)
    connector.line.width = PptxPt(1.5)
    return connector


def create_architecture_pptx():
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    GREEN = PptxRGB(0x10, 0xB9, 0x81)
    BLUE = PptxRGB(0x3B, 0x82, 0xF6)
    PURPLE = PptxRGB(0x8B, 0x5C, 0xF6)
    ORANGE = PptxRGB(0xF5, 0x9E, 0x0B)
    RED = PptxRGB(0xEF, 0x44, 0x44)
    DARK_CARD = PptxRGB(0x1E, 0x29, 0x3B)
    WHITE = PptxRGB(0xFF, 0xFF, 0xFF)
    GRAY = PptxRGB(0x9C, 0xA3, 0xAF)

    # ---- Slide 1: Title ----
    def slide1(slide, _):
        set_title(slide, "")
        add_text_box(slide, 1.5, 1.5, 10, 1, "ShieldLend", 48, WHITE, True, PP_ALIGN.CENTER)
        add_text_box(slide, 1.5, 2.5, 10, 0.6, "Privacy-Native BTC Lending Protocol on Starknet", 24, GREEN, False, PP_ALIGN.CENTER)
        add_text_box(slide, 1.5, 3.3, 10, 0.5, "Re{define} Hackathon  |  March 2026  |  by 0xhaz", 16, GRAY, False, PP_ALIGN.CENTER)

        # Key stats boxes
        stats = [
            ("16", "Cairo Contracts"),
            ("3,200+", "Lines of Code"),
            ("123", "Passing Tests"),
            ("3", "Live Markets"),
        ]
        for i, (num, label) in enumerate(stats):
            x = 2.0 + i * 2.5
            add_box(slide, x, 4.5, 2.0, 1.2, f"{num}\n{label}", DARK_CARD, 14, WHITE)

        add_text_box(slide, 1.5, 6.2, 10, 0.5, "Tracks:  Privacy  |  Bitcoin", 16, ORANGE, True, PP_ALIGN.CENTER)

    add_slide(prs, 6, "", slide1)

    # ---- Slide 2: Problem + Solution ----
    def slide2(slide, _):
        set_title(slide, "")
        add_text_box(slide, 0.8, 0.5, 12, 0.8, "The Problem & Our Solution", 32, WHITE, True, PP_ALIGN.LEFT)

        # Problem side
        add_box(slide, 0.8, 1.6, 5.5, 4.5,
                "THE PROBLEM\n\n"
                "Every DeFi lending transaction is\n"
                "fully transparent on-chain:\n\n"
                "- Deposit amounts visible to all\n"
                "- Borrow positions publicly tracked\n"
                "- Liquidation targets easily identified\n"
                "- Front-running and MEV extraction\n"
                "- No financial privacy for users",
                PptxRGB(0x7F, 0x1D, 0x1D), 14, WHITE)

        # Solution side
        add_box(slide, 7.0, 1.6, 5.5, 4.5,
                "OUR SOLUTION\n\n"
                "ShieldLend: Privacy as a first-class\n"
                "feature in BTC lending:\n\n"
                "- Pedersen commitments hide amounts\n"
                "- ZK proofs verify without revealing\n"
                "- Nullifiers prevent double-spends\n"
                "- Same pool, separate accounting\n"
                "- Native to Starknet (no external provers)",
                PptxRGB(0x06, 0x5F, 0x46), 14, WHITE)

    add_slide(prs, 6, "", slide2)

    # ---- Slide 3: Architecture Overview ----
    def slide3(slide, _):
        set_title(slide, "")
        add_text_box(slide, 0.8, 0.3, 12, 0.8, "Protocol Architecture", 32, WHITE, True, PP_ALIGN.LEFT)

        # Layer boxes
        # User Layer
        add_box(slide, 0.8, 1.4, 11.7, 0.8, "User Layer:    Next.js 16 Frontend  →  starknet-react Wallet  →  Starknet Sepolia", DARK_CARD, 13, GRAY)

        # Core Markets
        add_text_box(slide, 0.8, 2.4, 3, 0.4, "Core Markets", 14, GREEN, True)
        add_box(slide, 0.8, 2.8, 3.6, 0.6, "strkBTC / USDC  (75% LTV)", DARK_CARD, 12, WHITE)
        add_box(slide, 0.8, 3.5, 3.6, 0.6, "tBTC / USDC  (70% LTV)", DARK_CARD, 12, WHITE)
        add_box(slide, 0.8, 4.2, 3.6, 0.6, "strkBTC / tBTC eMode  (95% LTV)", DARK_CARD, 12, WHITE)

        # Supporting Contracts
        add_text_box(slide, 4.8, 2.4, 3, 0.4, "Supporting", 14, BLUE, True)
        add_box(slide, 4.8, 2.8, 3.6, 0.6, "SlToken  (deposit receipt)", DARK_CARD, 12, WHITE)
        add_box(slide, 4.8, 3.5, 3.6, 0.6, "DebtToken  (debt tracking)", DARK_CARD, 12, WHITE)
        add_box(slide, 4.8, 4.2, 3.6, 0.6, "InterestRateModel  (two-slope)", DARK_CARD, 12, WHITE)
        add_box(slide, 4.8, 4.9, 3.6, 0.6, "PragmaAdapter  (oracle)", DARK_CARD, 12, WHITE)

        # Privacy Engine
        add_text_box(slide, 8.8, 2.4, 3, 0.4, "Privacy Engine", 14, PURPLE, True)
        add_box(slide, 8.8, 2.8, 3.7, 0.6, "CommitmentStore  (Merkle tree)", DARK_CARD, 12, WHITE)
        add_box(slide, 8.8, 3.5, 3.7, 0.6, "NullifierRegistry  (double-spend)", DARK_CARD, 12, WHITE)
        add_box(slide, 8.8, 4.2, 3.7, 0.6, "ZkVerifier  (proof checking)", DARK_CARD, 12, WHITE)
        add_box(slide, 8.8, 4.9, 3.7, 0.6, "ShieldedAccount  (user state)", DARK_CARD, 12, WHITE)

        # Extra Layer
        add_text_box(slide, 0.8, 5.2, 4, 0.4, "Liquidation + Yield", 14, ORANGE, True)
        add_box(slide, 0.8, 5.6, 2.8, 0.6, "LiquidationEngine", DARK_CARD, 12, WHITE)
        add_box(slide, 3.8, 5.6, 2.8, 0.6, "DutchAuction", DARK_CARD, 12, WHITE)
        add_box(slide, 6.8, 5.6, 2.8, 0.6, "YieldTokenizer", DARK_CARD, 12, WHITE)
        add_box(slide, 9.8, 5.6, 2.7, 0.6, "FlashLoanVault", DARK_CARD, 12, WHITE)

        # Tech stack bar
        add_box(slide, 0.8, 6.5, 11.7, 0.7,
                "Cairo 2.14  |  Scarb 2.11.2  |  snforge 0.57.0  |  Next.js 16  |  React 19  |  Tailwind v4  |  Zustand v5",
                PptxRGB(0x11, 0x18, 0x27), 11, GRAY)

    add_slide(prs, 6, "", slide3)

    # ---- Slide 4: Privacy Flow ----
    def slide4(slide, _):
        set_title(slide, "")
        add_text_box(slide, 0.8, 0.3, 12, 0.8, "Shielded (Privacy) Lending Flow", 32, WHITE, True, PP_ALIGN.LEFT)
        add_text_box(slide, 0.8, 0.9, 12, 0.4, "Same LendingPool contract, different entrypoints — liquidity shared, accounting separate", 14, GRAY)

        # Step boxes — horizontal flow
        steps = [
            ("1. Client Prep", "Generate random\nblinding factor\n\nCompute Pedersen\ncommitment\n= Hash(amount, blind)", PURPLE),
            ("2. Shielded Deposit", "Submit commitment\n+ ZK proof to pool\n\nZkVerifier checks proof\n\nCommitmentStore inserts\ninto Merkle tree", GREEN),
            ("3. Save Note", "Store commitment +\nblinding + amount\nin local storage\n\n(Never goes on-chain)", DARK_CARD),
            ("4. Shielded Borrow", "Compute nullifier\nfrom secret\n\nZkVerifier checks proof\n\nNullifierRegistry\nprevents double-spend\n\nBorrow commitment\nadded to Merkle tree", BLUE),
        ]

        for i, (title, body, color) in enumerate(steps):
            x = 0.5 + i * 3.2
            add_box(slide, x, 1.6, 2.8, 4.8, f"{title}\n\n{body}", color, 12, WHITE)

        # Arrow hints
        for i in range(3):
            x = 3.3 + i * 3.2
            add_text_box(slide, x, 3.5, 0.5, 0.4, "→", 28, GRAY, True, PP_ALIGN.CENTER)

        # Bottom note
        add_box(slide, 0.5, 6.6, 12.3, 0.6,
                "Privacy Track:  Pedersen commitments  +  Merkle tree storage  +  Nullifier-based double-spend prevention  +  On-chain STARK verification",
                PptxRGB(0x1E, 0x1B, 0x4B), 12, PptxRGB(0xA7, 0x8B, 0xFA))

    add_slide(prs, 6, "", slide4)

    # ---- Slide 5: Lending Flow ----
    def slide5(slide, _):
        set_title(slide, "")
        add_text_box(slide, 0.8, 0.3, 12, 0.8, "Transparent Lending Flow", 32, WHITE, True, PP_ALIGN.LEFT)

        # 4 flow boxes
        flows = [
            ("Deposit", "User approves pool\n→ deposit(amount)\n→ Pool takes collateral\n→ Pool mints slTokens\n   to user (1:1)", PptxRGB(0x06, 0x5F, 0x46)),
            ("Borrow", "User calls borrow(amount)\n→ Pool checks LTV via Oracle\n→ collateral_value * LTV\n   >= debt_value\n→ Pool mints DebtToken\n→ Pool sends USDC to user", PptxRGB(0x0C, 0x4A, 0x6E)),
            ("Repay", "User calls repay(amount)\n→ Pool takes USDC from user\n→ Pool burns DebtToken\n→ Interest accrued via\n   borrow_index scaling", PptxRGB(0x92, 0x40, 0x0E)),
            ("Withdraw", "User calls withdraw(amount)\n→ Pool checks health factor\n→ Pool burns slTokens\n→ Pool returns collateral\n   to user", PptxRGB(0x1E, 0x1B, 0x4B)),
        ]

        for i, (title, body, color) in enumerate(flows):
            x = 0.5 + i * 3.2
            add_box(slide, x, 1.4, 2.8, 3.8, f"{title}\n\n{body}", color, 12, WHITE)

        # Interest rate model
        add_text_box(slide, 0.8, 5.5, 12, 0.4, "Interest Rate Model (Two-Slope)", 18, ORANGE, True)
        add_box(slide, 0.5, 6.0, 6.0, 1.2,
                "Below optimal utilization:\n  rate = base_rate + slope1 × (util / optimal)\n\n"
                "Above optimal utilization:\n  rate = base_rate + slope1 + slope2 × (excess / remaining)",
                DARK_CARD, 11, WHITE)
        add_box(slide, 6.8, 6.0, 6.0, 1.2,
                "Supply APY = borrow_rate × utilization × (1 - reserve_factor)\n\n"
                "Live rates:  strkBTC/USDC → 4.5% borrow / 2.0% supply\n"
                "             tBTC/USDC → 3.9% borrow / 1.3% supply",
                DARK_CARD, 11, GREEN)

    add_slide(prs, 6, "", slide5)

    # ---- Slide 6: Yield Tokenization + BTC Track ----
    def slide6(slide, _):
        set_title(slide, "")
        add_text_box(slide, 0.8, 0.3, 12, 0.8, "Yield Tokenization & BTC-Native Features", 32, WHITE, True, PP_ALIGN.LEFT)

        # Yield flow
        add_text_box(slide, 0.8, 1.2, 5, 0.4, "Yield Tokenization", 18, PURPLE, True)
        yield_steps = [
            ("Tokenize", "slTokens → PT + YT\n(June 30, 2026 maturity)"),
            ("Trade", "Fixed yield: sell YT, keep PT\nYield leverage: buy extra YT"),
            ("Redeem", "After maturity:\nPT → slTokens 1:1\nYT → accrued interest"),
        ]
        for i, (title, body) in enumerate(yield_steps):
            add_box(slide, 0.5 + i * 2.2, 1.8, 2.0, 2.5, f"{title}\n\n{body}", DARK_CARD, 11, WHITE)

        # BTC Track
        add_text_box(slide, 7.2, 1.2, 5, 0.4, "Bitcoin Track Features", 18, ORANGE, True)
        btc_features = [
            ("strkBTC", "Starknet-native\nwrapped BTC"),
            ("tBTC", "Threshold Network\nBTC bridge"),
            ("eMode", "95% LTV for\ncorrelated BTC pairs"),
        ]
        for i, (title, body) in enumerate(btc_features):
            add_box(slide, 7.0 + i * 2.2, 1.8, 2.0, 2.5, f"{title}\n\n{body}", DARK_CARD, 11, WHITE)

        # Additional features
        add_text_box(slide, 0.8, 4.6, 12, 0.4, "Additional Features", 18, GREEN, True)
        extras = [
            ("Flash Loans", "Uncollateralized atomic loans\n0.09% fee (9 bps)\nSame-tx borrow + repay"),
            ("Dutch Auction\nLiquidation", "Time-decaying price auctions\n2-5% liquidation bonus\nHealthy markets maintained"),
            ("Isolated Markets", "Each pool = separate risk\nNo cross-contamination\nIndependent rate models"),
            ("Full Frontend", "Next.js 16 + React 19\nWallet connection\nLive on-chain data"),
        ]
        for i, (title, body) in enumerate(extras):
            add_box(slide, 0.5 + i * 3.2, 5.1, 2.8, 2.0, f"{title}\n\n{body}", DARK_CARD, 11, WHITE)

    add_slide(prs, 6, "", slide6)

    # ---- Slide 7: Summary ----
    def slide7(slide, _):
        set_title(slide, "")
        add_text_box(slide, 1.5, 1.0, 10, 1, "ShieldLend", 48, WHITE, True, PP_ALIGN.CENTER)
        add_text_box(slide, 1.5, 2.0, 10, 0.6, "Privacy-Native BTC Lending on Starknet", 24, GREEN, False, PP_ALIGN.CENTER)

        # Summary points
        points = [
            "16 Cairo smart contracts  |  3,200+ lines  |  123 passing tests",
            "3 live markets on Starknet Sepolia with real on-chain data",
            "Privacy engine: Pedersen commitments + nullifiers + ZK verification",
            "Yield tokenization: PT/YT splitting with fixed maturity",
            "Flash loans, Dutch auction liquidations, eMode for BTC pairs",
            "Full Next.js 16 frontend with wallet integration",
        ]
        for i, point in enumerate(points):
            add_text_box(slide, 2.5, 3.0 + i * 0.5, 8, 0.5, f"  {point}", 15, GRAY if i % 2 == 0 else WHITE, False, PP_ALIGN.LEFT)

        add_text_box(slide, 1.5, 6.2, 10, 0.5, "github.com/0xhaz  |  0xhaz", 16, GRAY, False, PP_ALIGN.CENTER)
        add_text_box(slide, 1.5, 6.7, 10, 0.4, "[PASTE ARCHITECTURE DIAGRAM HERE IF NEEDED]", 12, PptxRGB(0x55, 0x55, 0x55), False, PP_ALIGN.CENTER)

    add_slide(prs, 6, "", slide7)

    path = os.path.join(OUT_DIR, "ShieldLend_Architecture.pptx")
    prs.save(path)
    print(f"Created: {path}")
    return path


if __name__ == "__main__":
    create_demo_docx()
    create_architecture_pptx()
    print("\nDone! Files in docs/")
